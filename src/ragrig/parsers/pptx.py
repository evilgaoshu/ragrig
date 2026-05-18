from __future__ import annotations

from hashlib import sha256
from pathlib import Path

from ragrig.parsers.base import ParseResult


class PptxParserError(ValueError):
    """Raised when a PPTX file cannot be parsed."""


class PptxParser:
    parser_name = "pptx"
    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def parse(self, path: Path) -> ParseResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise PptxParserError(
                "python-pptx is required for PPTX parsing: pip install python-pptx"
            ) from exc

        raw_bytes = path.read_bytes()
        try:
            prs = Presentation(path)
        except Exception as exc:
            raise PptxParserError(f"PPTX file could not be read: {path.name}") from exc

        slide_texts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            body_lines: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Identify title by placeholder type (idx 0 = title, idx 1 = body)
                ph = getattr(shape, "placeholder_format", None)
                if ph is not None and ph.idx == 0:
                    title = text
                else:
                    body_lines.append(text)

            header = f"[Slide {i}: {title}]" if title else f"[Slide {i}]"
            slide_text = header
            if body_lines:
                slide_text += "\n" + "\n".join(body_lines)
            slide_texts.append(slide_text)

        if not slide_texts:
            raise PptxParserError(f"PPTX file has no readable slides: {path.name}")

        extracted_text = "\n\n".join(slide_texts)
        return ParseResult(
            extracted_text=extracted_text,
            content_hash=sha256(raw_bytes).hexdigest(),
            mime_type=self.mime_type,
            parser_name=self.parser_name,
            metadata={
                "parser_id": "parser.pptx",
                "status": "success",
                "extension": path.suffix.lower(),
                "slide_count": len(prs.slides),
                "char_count": len(extracted_text),
            },
        )
